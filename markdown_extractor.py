#!/usr/bin/env python3
"""
Markdown Extractor v2.0 - Work Stream 1

Transform original source files into navigable Markdown working views.

Human Primacy:
- Source files are NEVER modified
- Markdown is a lossy working view
- Markdown may be regenerated (overwrites previous)
"""

import hashlib
import logging
from pathlib import Path
from typing import Dict, Any, List, Optional
from datetime import datetime


logger = logging.getLogger(__name__)


# Supported source file types
SUPPORTED_EXTENSIONS = {
    # Documents
    ".pdf", ".ppt", ".pptx", ".doc", ".docx",
    # Spreadsheets
    ".xls", ".xlsx",
    # Email
    ".eml", ".msg",
    # Text
    ".txt",
    # Images
    ".png", ".jpg", ".jpeg",
}


def discover_files(sources_path: Path) -> List[Path]:
    """Discover all supported source files recursively.

    Args:
        sources_path: Root directory to scan

    Returns:
        List of file paths
    """
    logger.info(f"[DISCOVER] Scanning directory: {sources_path}")

    files = []
    for file_path in sources_path.rglob("*"):
        # Skip directories
        if file_path.is_dir():
            continue

        # Check if extension is supported
        if file_path.suffix.lower() in SUPPORTED_EXTENSIONS:
            files.append(file_path)
            logger.debug(f"[DISCOVER] Found: {file_path}")
        else:
            logger.debug(f"[DISCOVER] Skipping unsupported: {file_path}")

    logger.info(f"[DISCOVER] Found {len(files)} supported files")
    return files


def generate_source_id(file_path: Path, sources_root: Path) -> str:
    """Generate unique source_id for a file.

    Uses path hash strategy (BQ-2 decision from plan).

    Format: src-{hash} where hash is first 12 chars of SHA256 of relative path

    Args:
        file_path: Absolute path to source file
        sources_root: Root directory for sources

    Returns:
        Unique source_id string
    """
    # Get relative path from sources root
    try:
        rel_path = file_path.relative_to(sources_root)
    except ValueError:
        # File is not under sources root, use absolute path
        rel_path = file_path

    # Hash the relative path
    path_str = str(rel_path)
    hash_obj = hashlib.sha256(path_str.encode('utf-8'))
    hash_hex = hash_obj.hexdigest()[:12]

    source_id = f"src-{hash_hex}"

    logger.debug(f"[SOURCE_ID] {file_path} -> {source_id}")
    return source_id


def extract_text_from_file(file_path: Path) -> str:
    """Extract text content from a source file.

    Args:
        file_path: Path to source file

    Returns:
        Extracted text content

    Raises:
        RuntimeError: If extraction fails
    """
    extension = file_path.suffix.lower()
    logger.info(f"[EXTRACT] Processing: {file_path} ({extension})")

    try:
        if extension == ".pdf":
            return extract_from_pdf(file_path)
        elif extension in [".ppt", ".pptx"]:
            return extract_from_powerpoint(file_path)
        elif extension in [".doc", ".docx"]:
            return extract_from_word(file_path)
        elif extension in [".xls", ".xlsx"]:
            return extract_from_excel(file_path)
        elif extension in [".eml", ".msg"]:
            return extract_from_email(file_path)
        elif extension == ".txt":
            return extract_from_text(file_path)
        elif extension in [".png", ".jpg", ".jpeg"]:
            return extract_from_image(file_path)
        else:
            raise RuntimeError(f"Unsupported file type: {extension}")

    except Exception as e:
        error_msg = f"Failed to extract from {file_path}: {e}"
        logger.error(f"[EXTRACT] {error_msg}")
        # Raise to allow caller to handle
        raise RuntimeError(error_msg) from e


def extract_from_pdf(file_path: Path) -> str:
    """Extract text from PDF.

    Args:
        file_path: Path to PDF file

    Returns:
        Extracted text
    """
    try:
        import PyPDF2
    except ImportError:
        raise RuntimeError("PyPDF2 is required for PDF extraction. Install: pip install PyPDF2")

    text_parts = []
    with open(file_path, 'rb') as f:
        reader = PyPDF2.PdfReader(f)
        num_pages = len(reader.pages)
        logger.debug(f"[EXTRACT_PDF] Pages: {num_pages}")

        for page_num, page in enumerate(reader.pages, 1):
            try:
                page_text = page.extract_text()
                if page_text:
                    text_parts.append(f"## Page {page_num}\n\n{page_text}\n")
            except Exception as e:
                logger.warning(f"[EXTRACT_PDF] Failed to extract page {page_num}: {e}")

    result = "\n".join(text_parts)
    logger.info(f"[EXTRACT_PDF] Extracted {len(result)} characters")
    return result


def extract_from_powerpoint(file_path: Path) -> str:
    """Extract text from PowerPoint.

    Args:
        file_path: Path to PowerPoint file

    Returns:
        Extracted text
    """
    try:
        from pptx import Presentation
    except ImportError:
        raise RuntimeError("python-pptx is required for PowerPoint extraction. Install: pip install python-pptx")

    text_parts = []
    prs = Presentation(file_path)

    for slide_num, slide in enumerate(prs.slides, 1):
        slide_text = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                slide_text.append(shape.text)

        if slide_text:
            text_parts.append(f"## Slide {slide_num}\n\n" + "\n".join(slide_text) + "\n")

    result = "\n".join(text_parts)
    logger.info(f"[EXTRACT_PPT] Extracted {len(result)} characters from {len(prs.slides)} slides")
    return result


def extract_from_word(file_path: Path) -> str:
    """Extract text from Word document.

    Args:
        file_path: Path to Word file

    Returns:
        Extracted text
    """
    try:
        import docx
    except ImportError:
        raise RuntimeError("python-docx is required for Word extraction. Install: pip install python-docx")

    doc = docx.Document(file_path)

    text_parts = []
    for para in doc.paragraphs:
        if para.text.strip():
            text_parts.append(para.text)

    result = "\n\n".join(text_parts)
    logger.info(f"[EXTRACT_DOCX] Extracted {len(result)} characters")
    return result


def extract_from_excel(file_path: Path) -> str:
    """Extract text from Excel spreadsheet.

    Args:
        file_path: Path to Excel file

    Returns:
        Extracted text
    """
    try:
        import openpyxl
    except ImportError:
        raise RuntimeError("openpyxl is required for Excel extraction. Install: pip install openpyxl")

    text_parts = []
    wb = openpyxl.load_workbook(file_path, data_only=True)

    for sheet_name in wb.sheetnames:
        sheet = wb[sheet_name]
        text_parts.append(f"## Sheet: {sheet_name}\n")

        for row in sheet.iter_rows(values_only=True):
            # Filter out None values and convert to strings
            row_text = " | ".join(str(cell) if cell is not None else "" for cell in row)
            if row_text.strip():
                text_parts.append(row_text)

        text_parts.append("")  # Blank line between sheets

    result = "\n".join(text_parts)
    logger.info(f"[EXTRACT_XLSX] Extracted {len(result)} characters from {len(wb.sheetnames)} sheets")
    return result


def extract_from_email(file_path: Path) -> str:
    """Extract text from email file.

    Args:
        file_path: Path to email file

    Returns:
        Extracted text
    """
    try:
        import email
        from email import policy
    except ImportError:
        raise RuntimeError("email module is required (built-in to Python)")

    with open(file_path, 'rb') as f:
        msg = email.message_from_binary_file(f, policy=policy.default)

    text_parts = []

    # Headers
    headers = [
        f"From: {msg.get('from', '')}",
        f"To: {msg.get('to', '')}",
        f"Subject: {msg.get('subject', '')}",
        f"Date: {msg.get('date', '')}",
    ]
    text_parts.append("\n".join(headers))
    text_parts.append("")

    # Body
    body = msg.get_body(preferencelist=('plain', 'html'))
    if body:
        content = body.get_content()
        text_parts.append(content)

    result = "\n".join(text_parts)
    logger.info(f"[EXTRACT_EMAIL] Extracted {len(result)} characters")
    return result


def extract_from_text(file_path: Path) -> str:
    """Extract text from text file.

    Args:
        file_path: Path to text file

    Returns:
        Extracted text
    """
    # Try different encodings
    encodings = ['utf-8', 'latin-1', 'cp1252']

    for encoding in encodings:
        try:
            with open(file_path, 'r', encoding=encoding) as f:
                result = f.read()
            logger.info(f"[EXTRACT_TXT] Extracted {len(result)} characters (encoding: {encoding})")
            return result
        except UnicodeDecodeError:
            continue

    # If all fail, read as binary and decode with errors='replace'
    with open(file_path, 'rb') as f:
        result = f.read().decode('utf-8', errors='replace')

    logger.info(f"[EXTRACT_TXT] Extracted {len(result)} characters (with replacement)")
    return result


def extract_from_image(file_path: Path) -> str:
    """Extract text from image using OCR.

    Note: This is a placeholder. In production, would use OCR.
    For now, returns metadata about the image.

    Args:
        file_path: Path to image file

    Returns:
        Extracted text (or metadata)
    """
    # Placeholder: OCR is not implemented for v2.0
    # Return metadata instead
    text = f"""# Image File

**Filename**: {file_path.name}
**Path**: {file_path}
**Size**: {file_path.stat().st_size} bytes

*Note: OCR extraction not implemented. This file contains an image that may have text content.*
"""

    logger.info(f"[EXTRACT_IMAGE] Generated metadata for {file_path}")
    return text


def generate_markdown(
    source_id: str,
    file_path: Path,
    sources_root: Path,
    extracted_text: str
) -> str:
    """Generate Markdown document with frontmatter.

    Args:
        source_id: Unique source ID
        file_path: Path to source file
        sources_root: Root directory for sources
        extracted_text: Extracted text content

    Returns:
        Complete Markdown document
    """
    # Get file stats
    stat = file_path.stat()
    rel_path = file_path.relative_to(sources_root)

    # Generate frontmatter
    frontmatter = f"""---
source_id: {source_id}
source_file: {file_path}
relative_path: {rel_path}
source_type: {file_path.suffix[1:]}  # Remove the dot
file_size: {stat.st_size}
extracted_at: {datetime.utcnow().isoformat()}Z
---

"""

    # Combine
    markdown = frontmatter + f"# {file_path.name}\n\n" + extracted_text

    return markdown


def extract_markdown(
    sources_path: Path,
    output_dir: Path,
    config: Dict[str, Any]
) -> None:
    """Extract Markdown from all source files.

    Main entry point for Work Stream 1.

    Args:
        sources_path: Path to source files directory
        output_dir: Path to output directory for Markdown files
        config: Configuration dictionary

    Raises:
        RuntimeError: If critical failures occur
    """
    logger.info(f"[EXTRACT_MARKDOWN] Starting operation...")
    logger.info(f"[EXTRACT_MARKDOWN] Sources: {sources_path}")
    logger.info(f"[EXTRACT_MARKDOWN] Output: {output_dir}")

    # Discover files
    files = discover_files(sources_path)

    if not files:
        logger.warning(f"[EXTRACT_MARKDOWN] No supported files found in {sources_path}")
        return

    # Create output directory
    output_dir.mkdir(parents=True, exist_ok=True)

    # Process each file
    success_count = 0
    fail_count = 0

    for file_path in files:
        try:
            # Generate source ID
            source_id = generate_source_id(file_path, sources_path)

            # Extract text
            extracted_text = extract_text_from_file(file_path)

            # Generate Markdown
            markdown = generate_markdown(
                source_id=source_id,
                file_path=file_path,
                sources_root=sources_path,
                extracted_text=extracted_text
            )

            # Write to output
            output_file = output_dir / f"{source_id}.md"
            with open(output_file, 'w', encoding='utf-8') as f:
                f.write(markdown)

            logger.info(f"[EXTRACT_MARKDOWN] Created: {output_file}")
            success_count += 1

        except Exception as e:
            logger.error(f"[EXTRACT_MARKDOWN] Failed to process {file_path}: {e}")
            fail_count += 1
            continue

    # Summary
    logger.info(f"[EXTRACT_MARKDOWN] Completed: {success_count} succeeded, {fail_count} failed")

    if fail_count > 0:
        logger.warning(f"[EXTRACT_MARKDOWN] {fail_count} files failed to process")

    if success_count == 0:
        raise RuntimeError("All files failed to process")
